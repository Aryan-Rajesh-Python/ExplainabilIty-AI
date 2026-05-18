# ============================================================
# UNIVERSAL EXPLAINABLE AI SYSTEM
# Full Multimodal XAI Platform
#
# Supports:
# - TXT
# - PDF
# - DOCX
# - PPTX
# - Images
# - Audio
# - Video
#
# Streamlit UI
# ============================================================

import os
import tempfile
import time
from collections import Counter
from concurrent.futures import ThreadPoolExecutor, TimeoutError as FuturesTimeout

import streamlit as st
import numpy as np
import nltk
import google.generativeai as genai

# ============================================================
# NLP / AI
# ============================================================

from transformers import pipeline
from sentence_transformers import SentenceTransformer

# ============================================================
# FILE PROCESSING
# ============================================================

import fitz
from docx import Document
from pptx import Presentation

# ============================================================
# IMAGE
# ============================================================

from PIL import Image

import torch
from torchvision import models, transforms
from torchvision.models import ResNet50_Weights

# ============================================================
# AUDIO
# ============================================================

import librosa
import soundfile as sf
import speech_recognition as sr
import whisper

@st.cache_resource
def load_whisper():
    return whisper.load_model("base")

# ============================================================
# VIDEO
# ============================================================

import cv2
from moviepy import VideoFileClip
import matplotlib.pyplot as plt

from xai.report import MechanisticReport, report_to_context
from xai.text import analyze_text as mechanistic_text
from xai.image import analyze_image as mechanistic_image
from xai.audio import analyze_audio as mechanistic_audio
from xai.video import analyze_video as mechanistic_video
from xai.document import (
    analyze_pdf_structure,
    analyze_docx_structure,
    analyze_pptx_structure,
)
from xai.discourse import analyze_discourse
from xai.audio_temporal import enrich_audio_report
from xai.video_temporal import enrich_video_temporal
from xai.pptx_discourse import analyze_pptx_discourse
from xai import viz as xai_viz
from xai.diffusion_trace import (
    generate_with_trace,
    load_diffusion_pipeline,
    trace_to_report,
)

# ============================================================
# DOWNLOAD NLTK
# ============================================================

@st.cache_resource
def download_nltk():
    nltk.download("punkt")
    nltk.download("punkt_tab")
    nltk.download("stopwords")

download_nltk()

# ============================================================
# STREAMLIT CONFIG
# ============================================================

st.set_page_config(
    page_title="Universal Explainable AI",
    layout="wide"
)

# ============================================================
# TITLE
# ============================================================

st.title("🧠 Universal Multimodal Explainable AI Framework")
st.caption(
    "Model-process-level XAI: feature attribution, inference traces, and generation pathways."
)

# ============================================================
# SIDEBAR — XAI OPTIONS
# ============================================================

with st.sidebar:
    st.header("XAI Engine")
    fast_mode = st.checkbox(
        "Fast mode (PDF / long documents)",
        value=True,
        help="Skips slow LIME/SHAP; uses keywords + one zero-shot pass (~30s vs 15+ min).",
    )
    use_shap = st.checkbox(
        "SHAP token attribution (text, slower)",
        value=False,
        help="Only applies when Fast mode is off. Can take many minutes on CPU.",
    )
    use_clip = st.checkbox(
        "CLIP prompt–region heatmap (images)",
        value=True,
        help="Cross-modal patch alignment vs your prompt.",
    )
    st.divider()
    st.header("Diffusion Lab")
    st.caption(
        "Generate in-app to capture denoising-step traces (GPU recommended)."
    )
    run_diffusion_lab = st.checkbox("Enable diffusion generation", value=False)
    diffusion_steps = st.slider("Denoising steps", 10, 30, 20)
    diffusion_model = st.selectbox(
        "Model",
        [
            "runwayml/stable-diffusion-v1-5",
            "stabilityai/stable-diffusion-2-1-base",
        ],
    )
    st.divider()
    app_mode = st.radio(
        "Application mode",
        ["Multimodal XAI", "Explainable Travel Agent"],
        help="Travel Agent: plan trips + explain why the agent chose that plan.",
    )

# ============================================================
# GEMINI CONFIG
# ============================================================

GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")
gemini_configured = False

if GEMINI_API_KEY:
    genai.configure(api_key=GEMINI_API_KEY)
    gemini_configured = True
else:
    st.error(
        "Gemini API key is not configured. Set GEMINI_API_KEY in your environment."
    )

# ============================================================
# LIST AVAILABLE GEMINI MODELS
# ============================================================

available_models = []

if gemini_configured:

    try:

        for m in genai.list_models():

            if "generateContent" in m.supported_generation_methods:

                available_models.append(m.name)

    except Exception as e:

        st.error(
            f"Could not load Gemini models: {e}"
        )

else:

    st.warning(
        "Gemini is not configured; model listing is disabled."
    )

# ============================================================
# MODEL SELECTOR
# ============================================================

if available_models:

    selected_model = st.selectbox(
        "Choose Gemini Model",
        available_models
    )

    model = genai.GenerativeModel(
        selected_model
    )

else:

    selected_model = None
    model = None
    st.warning(
        "No Gemini model available. Gemini reasoning is disabled until API key and model access are configured."
    )

# ============================================================
# LOAD MODELS
# ============================================================

@st.cache_resource
def load_models():

    sentiment_model = pipeline(
        "sentiment-analysis"
    )

    zero_shot_model = pipeline(
        "zero-shot-classification",
        model="facebook/bart-large-mnli"
    )

    embedding_model = SentenceTransformer(
        "all-MiniLM-L6-v2"
    )

    # ============================================================
    # LOAD RESNET50 + IMAGENET LABELS
    # ============================================================

    weights = ResNet50_Weights.DEFAULT

    image_model = models.resnet50(
        weights=weights
    )

    image_model.eval()

    imagenet_labels = weights.meta["categories"]

    return (
        sentiment_model,
        zero_shot_model,
        embedding_model,
        image_model,
        imagenet_labels
    )

(
    sentiment_model,
    zero_shot_model,
    embedding_model,
    image_model,
    imagenet_labels
) = load_models()


@st.cache_resource
def load_clip():
    from transformers import CLIPModel, CLIPProcessor

    model_id = "openai/clip-vit-base-patch32"
    processor = CLIPProcessor.from_pretrained(model_id)
    model = CLIPModel.from_pretrained(model_id)
    model.eval()
    return model, processor


@st.cache_resource
def get_diffusion_pipe(model_id):
    return load_diffusion_pipeline(model_id)


def get_clip():
    return load_clip()


def reset_xai_session():
    """Clear per-run state so explanations cannot leak across uploads."""
    for key in list(st.session_state.keys()):
        if key.startswith("xai_"):
            del st.session_state[key]
    st.session_state["xai_run"] = id(object())


# ============================================================
# FILE HELPERS
# ============================================================

PDF_MAX_CHARS = 30_000


def safe_unlink(path, retries=8, delay=0.25):
    """Windows-safe temp file removal (handles open handles from st.video / moviepy)."""
    if not path or not os.path.exists(path):
        return
    for attempt in range(retries):
        try:
            os.unlink(path)
            return
        except PermissionError:
            if attempt < retries - 1:
                time.sleep(delay)
        except OSError:
            return


def extract_pdf_text(path, max_chars=PDF_MAX_CHARS):

    doc = fitz.open(path)

    text = ""

    for page in doc:
        text += page.get_text()
        if len(text) >= max_chars:
            text = text[:max_chars]
            break

    doc.close()
    return text


def extract_docx_text(path):

    doc = Document(path)

    text = []

    for para in doc.paragraphs:
        text.append(para.text)

    return "\n".join(text)


def extract_pptx_text(path):

    prs = Presentation(path)

    text = []

    for slide in prs.slides:

        for shape in slide.shapes:

            if hasattr(shape, "text"):

                text.append(shape.text)

    return "\n".join(text)

# ============================================================
# TEXT EXPLAINABILITY
# ============================================================

def explain_text(text):

    explanation = {}

    words = text.split()

    explanation["Word Count"] = len(words)

    explanation["Character Count"] = len(text)

    # ========================================================
    # SENTIMENT
    # ========================================================

    try:

        sentiment = sentiment_model(
            text[:512]
        )[0]

        explanation["Sentiment"] = sentiment

    except Exception as e:

        explanation["Sentiment"] = (
            f"Could not analyze sentiment: {str(e)}"
        )

    # ========================================================
    # KEYWORDS
    # ========================================================

    tokens = nltk.word_tokenize(
        text.lower()
    )

    stopwords = set(
        nltk.corpus.stopwords.words("english")
    )

    filtered = [

        word for word in tokens

        if word.isalpha()
        and word not in stopwords

    ]

    common = Counter(filtered).most_common(10)

    explanation["Top Keywords"] = common

    # ========================================================
    # SEMANTIC EMBEDDING
    # ========================================================

    embedding = embedding_model.encode([text[:8000]])

    explanation["Embedding Shape"] = (
        embedding.shape
    )

    # ========================================================
    # SEMANTIC STYLE ANALYSIS
    # ========================================================

    candidate_labels = [

        "summarization",
        "analytical writing",
        "recommendation",
        "academic writing",
        "creative writing",
        "educational content",
        "technical explanation",
        "storytelling"

    ]

    semantic_result = zero_shot_model(
        text[:2000],
        candidate_labels
    )

    reasoning = []

    for label, score in zip(
        semantic_result["labels"],
        semantic_result["scores"]
    ):

        if score > 0.45:

            reasoning.append(
                f"Strong semantic alignment detected with '{label}' patterns (confidence: {round(score, 2)})."
            )

    if len(words) > 300:

        reasoning.append(
            "Long-form structure suggests the AI prioritized depth, elaboration, and informational completeness."
        )

    explanation["Semantic Categories"] = list(zip(
        semantic_result["labels"],
        semantic_result["scores"]
    ))

    explanation["Reasoning"] = reasoning

    return explanation


def display_mechanistic_report(report: MechanisticReport):
    st.subheader("🔬 Mechanistic XAI Trace")
    sections = [
        ("Input Semantics", report.input_semantics),
        ("Feature Attribution", report.feature_attribution),
        ("Internal Representation", report.internal_representation),
        ("Generation Pathway", report.generation_pathway),
        ("Output Alignment", report.output_alignment),
    ]
    for title, items in sections:
        with st.expander(title, expanded=title == "Feature Attribution"):
            if items:
                for item in items:
                    st.markdown(f"- {item}")
            else:
                st.markdown("_No attributions in this section._")

    artifacts = report.artifacts
    if "gradcam" in artifacts:
        fig = xai_viz.plot_gradcam_overlay(
            artifacts.get("image_pil"),
            artifacts["gradcam"],
        )
        st.pyplot(fig)
        plt.close(fig)
    if "mel_spectrogram" in artifacts:
        fig = xai_viz.plot_mel_spectrogram(artifacts["mel_spectrogram"])
        st.pyplot(fig)
        plt.close(fig)
    if "zero_shot" in artifacts:
        labels, scores = zip(*artifacts["zero_shot"])
        fig = xai_viz.plot_zero_shot_bars(list(labels), list(scores))
        st.pyplot(fig)
        plt.close(fig)
    if "clip_heatmap" in artifacts and artifacts.get("image_pil"):
        fig = xai_viz.plot_clip_heatmap(
            artifacts["image_pil"],
            artifacts["clip_heatmap"],
        )
        st.pyplot(fig)
        plt.close(fig)
    if "token_attribution" in artifacts:
        fig = xai_viz.plot_token_attribution(
            artifacts["token_attribution"],
            title="Token attribution (SHAP / LIME)",
        )
        if fig:
            st.pyplot(fig)
            plt.close(fig)
    if "confidence_trajectory" in artifacts:
        fig = xai_viz.plot_confidence_trajectory(
            artifacts["confidence_trajectory"]
        )
        if fig:
            st.pyplot(fig)
            plt.close(fig)
    if "diffusion_steps" in artifacts:
        fig = xai_viz.plot_diffusion_steps(artifacts["diffusion_steps"])
        if fig:
            st.pyplot(fig)
            plt.close(fig)
    if artifacts.get("generated_image") is not None:
        st.image(
            artifacts["generated_image"],
            caption="Diffusion-generated output",
            use_container_width=True,
        )
    if artifacts.get("agent_decision_steps"):
        with st.expander("Agent decision steps (tool calls)", expanded=True):
            for step in artifacts["agent_decision_steps"]:
                st.markdown(
                    f"**Step {step['step_id']}:** `{step['tool_name']}` "
                    f"({step['duration_ms']:.0f} ms)"
                )
                st.caption(step["tool_description"])
                st.markdown(f"- **Input:** {step['input_summary']}")
                st.markdown(f"- **Output:** {step['output_summary']}")
                st.markdown(f"- **Why:** {step['rationale']}")
    if artifacts.get("trace_sources"):
        st.caption(
            "Computation sources: " + ", ".join(artifacts["trace_sources"])
        )


def cleanup_explanation_context(
    explanation,
    content_type="text"
):

    cleaned = []

    # ========================================================
    # TEXT / DOCUMENTS
    # ========================================================

    if content_type in ["text", "document"]:

        word_count = explanation.get("Word Count", 0)

        sentiment = explanation.get("Sentiment", {})
        sentiment_label = sentiment.get("label", "")

        keywords = explanation.get(
            "Top Keywords",
            []
        )

        reasoning = explanation.get(
            "Reasoning",
            []
        )

        if word_count > 300:

            cleaned.append(
                "The output prioritizes depth and detailed explanation."
            )

        elif word_count > 100:

            cleaned.append(
                "The response balances clarity with moderate detail."
            )

        else:

            cleaned.append(
                "The output focuses on concise communication."
            )

        if sentiment_label == "POSITIVE":

            cleaned.append(
                "The wording creates a positive and engaging tone."
            )

        elif sentiment_label == "NEGATIVE":

            cleaned.append(
                "The language introduces stronger emotional intensity."
            )

        if keywords:

            important_words = [

                word for word, _ in keywords[:5]

            ]

            cleaned.append(
                f"The response strongly emphasizes themes related to: {', '.join(important_words)}."
            )

        cleaned.extend(reasoning)

    # ========================================================
    # IMAGE
    # ========================================================

    elif content_type == "image":

        extracted_text = explanation.get(
            "Extracted Text",
            ""
        )

        image_type = explanation.get(
            "Detected Image Type",
            ""
        )

        if (
            extracted_text
            and extracted_text.strip() != ""
            and "ocr failed" not in extracted_text.lower()
            and "no readable text" not in extracted_text.lower()
        ):

            cleaned.append(
                "The image includes readable content that influences its overall presentation and structure."
            )

        if (
            image_type
            == "Natural Image / Photograph"
        ):

            cleaned.append(
                "The visual style focuses on realism, subject clarity, and cinematic composition."
            )

        elif (
            image_type
            == "Document / Presentation"
        ):

            cleaned.append(
                "The image emphasizes structured information and organized visual layout."
            )

        else:

            cleaned.append(
                "The image combines visual aesthetics with informational structure."
            )

        cleaned.append(
            "Lighting, composition, and visual emphasis strongly influence the final appearance."
        )

    # ========================================================
    # AUDIO
    # ========================================================

    elif content_type == "audio":

        transcript = explanation.get(
            "Transcript",
            ""
        )

        tempo = explanation.get(
            "Tempo",
            0
        )

        if transcript and not transcript.lower().startswith(
            "could not"
        ):

            cleaned.append(
                "Spoken language and vocal delivery strongly shape the generated result."
            )

        if tempo > 120:

            cleaned.append(
                "The pacing feels energetic and fast-moving."
            )

        else:

            cleaned.append(
                "The pacing feels calmer and more controlled."
            )

    # ========================================================
    # VIDEO
    # ========================================================

    elif content_type == "video":

        duration = explanation.get(
            "Duration",
            0
        )

        motion = explanation.get(
            "Average Motion",
            0
        )

        if motion > 25:

            cleaned.append(
                "The video emphasizes movement and dynamic visual transitions."
            )

        else:

            cleaned.append(
                "The video maintains a steadier and more controlled visual pace."
            )

        if duration > 60:

            cleaned.append(
                "The structure feels more narrative and progression-focused."
            )

        cleaned.append(
            "Visual pacing, scene composition, and sequencing shape the overall presentation."
        )

    # ========================================================
    # FINAL CLEANUP
    # ========================================================

    cleaned_text = "\n".join(cleaned)

    banned_terms = [

        "OCR",
        "confidence",
        "embedding",
        "pipeline",
        "classification",
        "detected",
        "prediction",
        "analysis failed",
        "transcription failed",
        "natural score",
        "document score"

    ]

    for term in banned_terms:

        cleaned_text = cleaned_text.replace(
            term,
            ""
        )

    return cleaned_text.strip()

# ============================================================
# GEMINI HUMAN EXPLAINABILITY
# ============================================================

def generate_human_explanation(
    user_prompt,
    context,
    image=None
):

    prompt = f"""
    You are a research-grade explainable AI narrator.

    Explain HOW the generation pipeline likely worked internally — not merely what appears in the output.

    You receive:
    1. The original user prompt
    2. A mechanistic XAI trace (feature attribution, inference traces, generation pathway)

    Use mechanistic XAI language:
    - inference trace (not "reasoning")
    - attributed / semantically prioritized (not "detected")
    - causally interpreted (not "analyzed")
    - emerged through (not "generated because")
    - latent aesthetic representation (not "style")

    Structure your narrative across these five layers:
    1. Input Semantics — strongest prompt/file semantics
    2. Feature Attribution — what features contributed most
    3. Internal Representation — how content was interpreted internally
    4. Generation Pathway — how the output evolved step-by-step
    5. Final Output Alignment — why the final output matched the prompt

    Be specific to the trace. Avoid generic summaries.

    CRITICAL causal rules:
    - Only use facts present in the mechanistic trace. Do not invent pipeline stages.
    - For uploaded audio/video: describe as "observed output artifact", NEVER as
      "reference audio", "reference conditioning", or "example speech input".
    - Cite computation sources when the trace includes them.
    - Do not fabricate exact percentages unless they appear in the trace.

    USER PROMPT:
    {user_prompt}

    MECHANISTIC XAI TRACE:
    {context}

    Write a cohesive, academic-quality mechanistic explanation.
    """

    if model is None:

        return "Gemini is not configured. Set GEMINI_API_KEY and ensure model access is available."

    def _call_gemini():
        if image:
            return model.generate_content([prompt, image])
        return model.generate_content(prompt)

    try:
        with ThreadPoolExecutor(max_workers=1) as pool:
            future = pool.submit(_call_gemini)
            response = future.result(timeout=90)
        return response.text

    except FuturesTimeout:
        return (
            "Gemini narrative timed out after 90s. "
            "The mechanistic XAI trace above is complete."
        )
    except Exception:
        return (
            "The explanation could not be generated "
            "for this content."
        )
    
# ============================================================
# OCR
# ============================================================

import pytesseract
pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"

# ============================================================
# ADVANCED IMAGE TYPE INFERENCE
# ============================================================

def infer_image_type(
    image,
    predictions,
    extracted_text
):

    # ========================================================
    # OCR TEXT DENSITY
    # ========================================================

    text_length = len(
        extracted_text.strip()
    )

    # ========================================================
    # VISUAL CONFIDENCE
    # ========================================================

    top_confidence = max(

        pred["confidence"]

        for pred in predictions

    )

    # ========================================================
    # OBJECT DIVERSITY
    # ========================================================

    unique_objects = len(set(

        pred["label"]

        for pred in predictions

    ))

    # ========================================================
    # IMAGE STATISTICS
    # ========================================================

    img_array = np.array(image)

    grayscale = cv2.cvtColor(
        img_array,
        cv2.COLOR_RGB2GRAY
    )

    # ========================================================
    # EDGE DENSITY
    # ========================================================

    edges = cv2.Canny(
        grayscale,
        100,
        200
    )

    edge_density = np.mean(edges)

    # ========================================================
    # TEXT COVERAGE ESTIMATION
    # ========================================================

    text_density = text_length / (
        image.size[0] * image.size[1]
    )

    # ========================================================
    # NATURAL IMAGE SCORE
    # ========================================================

    natural_score = (

        (top_confidence * 0.5)

        +

        (unique_objects * 0.1)

        +

        (edge_density / 255 * 0.4)

    )

    # ========================================================
    # DOCUMENT SCORE
    # ========================================================

    document_score = (

        (text_density * 1000 * 0.7)

        +

        ((1 - unique_objects / 10) * 0.3)

    )

    # ========================================================
    # FINAL DECISION
    # ========================================================

    if natural_score > document_score:

        return (
            "Natural Image / Photograph",
            natural_score,
            document_score
        )

    elif text_length > 300:

        return (
            "Document / Presentation",
            natural_score,
            document_score
        )

    else:

        return (
            "Mixed / Ambiguous Content",
            natural_score,
            document_score
        )
    
# ============================================================
# IMPROVED IMAGE EXPLAINABILITY
# ============================================================

def explain_image(image):

    explanation = {}

    # ========================================================
    # OCR TEXT EXTRACTION
    # ========================================================

    try:

        extracted_text = pytesseract.image_to_string(
            image
        )

    except Exception as e:

        extracted_text = f"OCR failed: {str(e)}"

    # ========================================================
    # EMPTY OCR FALLBACK
    # ========================================================

    if extracted_text.strip() == "":

        extracted_text = (
            "No readable text detected."
        )

    explanation["Extracted Text"] = (
        extracted_text[:3000]
    )

    # ========================================================
    # BASIC IMAGE METADATA
    # ========================================================

    explanation["Image Size"] = image.size

    explanation["Image Mode"] = image.mode

    # ========================================================
    # VISUAL ANALYSIS
    # ========================================================

    transform = transforms.Compose([

        transforms.Resize((224, 224)),

        transforms.ToTensor()

    ])

    img_tensor = transform(
        image
    ).unsqueeze(0)

    with torch.no_grad():

        outputs = image_model(img_tensor)

    probabilities = torch.nn.functional.softmax(
        outputs[0],
        dim=0
    )

    top5_prob, top5_catid = torch.topk(
        probabilities,
        5
    )

    # ========================================================
    # TOP PREDICTIONS
    # ========================================================

    predictions = []

    for i in range(top5_prob.size(0)):

        class_id = int(top5_catid[i])

        label = imagenet_labels[class_id]

        confidence = round(
            float(top5_prob[i]),
            4
        )

        predictions.append({

            "label": label,
            "confidence": confidence

        })

    explanation["Top Predictions"] = (
        predictions
    )

    # ========================================================
    # ADVANCED IMAGE TYPE INFERENCE
    # ========================================================

    (
        image_type,
        natural_score,
        document_score
    ) = infer_image_type(

        image,
        predictions,
        extracted_text

    )

    explanation["Detected Image Type"] = (
        image_type
    )

    explanation["Natural Score"] = round(
        natural_score,
        4
    )

    explanation["Document Score"] = round(
        document_score,
        4
    )

    # ========================================================
    # REAL REASONING
    # ========================================================

    reasoning = []

    if image_type == (
        "Document / Presentation"
    ):

        reasoning.append(
            "The image appears to be an educational or presentation slide containing structured textual information."
        )

        reasoning.append(
            "Large typography and centered formatting strongly influenced document-style interpretation."
        )

    if len(extracted_text) > 200:

        reasoning.append(
            "OCR detected substantial readable text, indicating semantic information is more important than object recognition."
        )

    if image.mode == "RGB":

        reasoning.append(
            "Color composition and contrast influenced visual segmentation."
        )

    reasoning.append(
        "Visual hierarchy and text prominence influenced semantic interpretation."
    )

    explanation["Reasoning"] = reasoning

    return explanation

# ============================================================
# AUDIO EXPLAINABILITY
# ============================================================

def explain_audio(audio_path):

    explanation = {}

    # ========================================================
    # LOAD AUDIO
    # ========================================================

    y, sr_rate = librosa.load(audio_path)

    duration = librosa.get_duration(
        y=y,
        sr=sr_rate
    )

    explanation["Duration"] = (
        round(duration, 2)
    )

    explanation["Sample Rate"] = sr_rate

    # ========================================================
    # ENERGY
    # ========================================================

    rms = librosa.feature.rms(y=y)[0]

    explanation["Average Energy"] = float(
        np.mean(rms)
    )

    # ========================================================
    # TEMPO
    # ========================================================

    tempo, _ = librosa.beat.beat_track(
        y=y,
        sr=sr_rate
    )

    # Fix for newer librosa versions

    if isinstance(tempo, np.ndarray):
        tempo = tempo[0]

    explanation["Tempo"] = float(tempo)

    # ========================================================
    # PITCH
    # ========================================================

    pitches, magnitudes = librosa.piptrack(
        y=y,
        sr=sr_rate
    )

    pitch_values = pitches[
        magnitudes > np.median(magnitudes)
    ]

    avg_pitch = (
        float(np.mean(pitch_values))
        if len(pitch_values) > 0
        else 0.0
    )

    explanation["Average Pitch"] = avg_pitch

    # ========================================================
    # TRANSCRIPTION
    # ========================================================

    transcript = ""

    try:

        whisper_model = load_whisper()

        result = whisper_model.transcribe(
            audio_path
        )

        transcript = result["text"]
        explanation["Whisper Segments"] = result.get("segments", [])

    except Exception as e:

        transcript = (
            f"Could not transcribe audio: {str(e)}"
        )

    explanation["Transcript"] = transcript

    # ========================================================
    # REASONING
    # ========================================================

    reasoning = []

    if tempo > 120:

        reasoning.append(
            "Fast tempo suggests energetic or intense generation."
        )

    if avg_pitch > 200:

        reasoning.append(
            "Higher pitch patterns indicate emotional emphasis."
        )

    if np.mean(rms) < 0.01:

        reasoning.append(
            "Low energy suggests calm or soft delivery."
        )

    if transcript and not transcript.lower().startswith(
        "could not transcribe audio"
    ):

        reasoning.append(
            "Speech semantics heavily influenced interpretation."
        )

    explanation["Reasoning"] = reasoning

    return explanation

# ============================================================
# VIDEO EXPLAINABILITY
# ============================================================

def explain_video(video_path):

    explanation = {}

    cap = cv2.VideoCapture(video_path)

    frame_count = int(
        cap.get(cv2.CAP_PROP_FRAME_COUNT)
    )

    fps = cap.get(cv2.CAP_PROP_FPS)

    duration = frame_count / fps

    explanation["Frame Count"] = frame_count

    explanation["FPS"] = fps

    explanation["Duration"] = round(duration, 2)

    sampled_frames = []

    current_frame = 0

    while True:

        ret, frame = cap.read()

        if not ret:
            break

        if current_frame % int(fps * 2) == 0:

            frame_rgb = cv2.cvtColor(
                frame,
                cv2.COLOR_BGR2RGB
            )

            sampled_frames.append(
                frame_rgb
            )

        current_frame += 1

    cap.release()

    explanation["Sampled Frames"] = len(
        sampled_frames
    )

    # ========================================================
    # MOTION ANALYSIS
    # ========================================================

    motion_scores = []

    for i in range(1, len(sampled_frames)):

        prev = cv2.cvtColor(
            sampled_frames[i - 1],
            cv2.COLOR_RGB2GRAY
        )

        curr = cv2.cvtColor(
            sampled_frames[i],
            cv2.COLOR_RGB2GRAY
        )

        diff = cv2.absdiff(prev, curr)

        motion_scores.append(
            np.mean(diff)
        )

    avg_motion = (
        np.mean(motion_scores)
        if len(motion_scores) > 0
        else 0
    )

    explanation["Average Motion"] = float(
        avg_motion
    )

    # ========================================================
    # AUDIO EXTRACTION
    # ========================================================

    transcript = "No transcript"
    whisper_segments = []
    clip = None
    audio_path = None

    try:
        clip = VideoFileClip(video_path)
        if clip.audio is None:
            raise ValueError("Video has no audio track")

        with tempfile.NamedTemporaryFile(
            suffix=".wav",
            delete=False,
        ) as audio_tmp:
            audio_path = audio_tmp.name

        clip.audio.write_audiofile(audio_path, logger=None)

        whisper_model = load_whisper()
        result = whisper_model.transcribe(audio_path)
        transcript = result["text"]
        whisper_segments = result.get("segments", [])

    except Exception:
        transcript = "Could not transcribe video"

    finally:
        if clip is not None:
            try:
                clip.close()
            except Exception:
                pass
        safe_unlink(audio_path)

    explanation["Transcript"] = transcript
    explanation["Whisper Segments"] = whisper_segments

    # ========================================================
    # REASONING
    # ========================================================

    reasoning = []

    if avg_motion > 25:

        reasoning.append(
            "High motion suggests dynamic scene generation."
        )

    if duration > 60:

        reasoning.append(
            "Long duration indicates narrative-driven structure."
        )

    if transcript != "Could not transcribe video":

        reasoning.append(
            "Speech/audio content heavily influenced semantic understanding."
        )

    explanation["Reasoning"] = reasoning
    explanation["motion_scores"] = motion_scores

    return explanation


def run_mechanistic_narrative(
    report,
    user_prompt,
    image=None,
    title="🧠 Mechanistic Inference Narrative",
    extra_context: str = "",
):
    display_mechanistic_report(report)
    context = report_to_context(report)
    if extra_context:
        context = context + "\n\n" + extra_context
    narrative = generate_human_explanation(
        user_prompt,
        context,
        image=image,
    )
    st.subheader(title)
    st.markdown(narrative)


# ============================================================
# EXPLAINABLE TRAVEL AGENT
# ============================================================

try:
    _app_mode = app_mode
except NameError:
    _app_mode = "Multimodal XAI"

if _app_mode == "Explainable Travel Agent":
    from travel_ui import render_travel_agent

    render_travel_agent(
        model,
        sentiment_model,
        zero_shot_model,
        embedding_model,
        fast_mode,
        run_mechanistic_narrative,
        reset_xai_session,
    )
    st.stop()


# ============================================================
# FILE UPLOADER
# ============================================================

user_prompt = st.text_area(
    "Enter the Original Prompt",
    height=200
)

uploaded_file = st.file_uploader(

    "Upload AI-generated Content",

    type=[

        "txt",
        "pdf",
        "docx",
        "pptx",

        "png",
        "jpg",
        "jpeg",

        "mp3",
        "wav",

        "mp4",
        "mov",
        "avi"

    ]
)

analyze_button = st.button(
    "Run Mechanistic XAI Analysis"
)

diffusion_button = (
    run_diffusion_lab
    and st.button("Generate image + diffusion trace")
)

# ============================================================
# DIFFUSION LAB (no upload required)
# ============================================================

if diffusion_button:

    reset_xai_session()

    if not user_prompt.strip():
        st.warning("Enter a prompt for diffusion generation.")
    else:
        with st.spinner(
            "Running Stable Diffusion with step-level latent tracing…"
        ):
            try:
                pipe_bundle = get_diffusion_pipe(diffusion_model)
                trace, _ = generate_with_trace(
                    user_prompt,
                    pipe_bundle=pipe_bundle,
                    num_inference_steps=diffusion_steps,
                    model_id=diffusion_model,
                )
                diff_report = trace_to_report(trace)
                st.image(
                    trace.image,
                    caption="Generated image",
                    use_container_width=True,
                )
                clip_m, clip_p = (get_clip() if use_clip else (None, None))
                img_report = mechanistic_image(
                    trace.image,
                    image_model,
                    imagenet_labels,
                    embedding_model,
                    user_prompt,
                    use_clip=use_clip,
                    clip_model=clip_m,
                    clip_processor=clip_p,
                )
                img_report.artifacts["image_pil"] = trace.image
                diff_report.extend(img_report)
                run_mechanistic_narrative(
                    diff_report,
                    user_prompt,
                    image=trace.image,
                    title="🧠 Diffusion Mechanistic Inference Narrative",
                )
            except Exception as e:
                st.error(
                    f"Diffusion generation failed: {e}. "
                    "Ensure diffusers/torch are installed; GPU recommended."
                )

# ============================================================
# PROCESS FILE
# ============================================================

if user_prompt.strip() and uploaded_file and analyze_button:

    reset_xai_session()

    # Security: File size limit
    if uploaded_file.size > 50 * 1024 * 1024:  # 50MB
        st.error("File too large. Maximum size is 50MB.")
        st.stop()

    suffix = uploaded_file.name.split(".")[-1].lower()
    file_bytes = uploaded_file.getvalue()

    with tempfile.NamedTemporaryFile(
        delete=False,
        suffix=f".{suffix}",
    ) as tmp:
        tmp.write(file_bytes)
        temp_path = tmp.name

    # Clean up temp file after processing
    try:

        st.success(
            "File uploaded successfully!"
        )

        # ========================================================
        # TXT
        # ========================================================

        if suffix == "txt":

            with open(
                temp_path,
                "r",
                encoding="utf-8"
            ) as f:

                text = f.read()

            st.subheader("📄 Text Content")

            st.text_area(
                "Extracted Text",
                text,
                height=300
            )

            with st.status("Running mechanistic XAI…", expanded=True) as status:
                status.write("Text routing & attribution…")
                report = mechanistic_text(
                    text,
                    sentiment_model,
                    zero_shot_model,
                    embedding_model,
                    user_prompt,
                    use_shap=use_shap and not fast_mode,
                    fast_mode=fast_mode,
                )
                report.extend(
                    analyze_discourse(
                        text,
                        user_prompt,
                        "txt",
                        embedding_model,
                        zero_shot_model,
                    )
                )
                status.write("Generating narrative…")
            run_mechanistic_narrative(report, user_prompt)

        # ========================================================
        # PDF
        # ========================================================

        elif suffix == "pdf":

            with st.status("Analyzing PDF…", expanded=True) as status:
                status.write("Extracting text…")
                text = extract_pdf_text(temp_path)
                if len(text) >= PDF_MAX_CHARS:
                    st.info(
                        f"PDF truncated to first {PDF_MAX_CHARS:,} characters for speed."
                    )

                st.subheader("📄 PDF Content")
                st.text_area(
                    "Extracted PDF Text",
                    text[:5000],
                    height=300,
                )

                status.write(
                    "Mechanistic trace (fast mode)…"
                    if fast_mode
                    else "Mechanistic trace (full LIME — may take several minutes)…"
                )
                report = mechanistic_text(
                    text,
                    sentiment_model,
                    zero_shot_model,
                    embedding_model,
                    user_prompt,
                    use_shap=use_shap and not fast_mode,
                    fast_mode=fast_mode,
                )
                report.extend(analyze_pdf_structure(text[:5000]))
                report.extend(
                    analyze_discourse(
                        text,
                        user_prompt,
                        "pdf",
                        embedding_model,
                        zero_shot_model,
                    )
                )
                status.write("Generating narrative…")
            run_mechanistic_narrative(report, user_prompt)

        # ========================================================
        # DOCX
        # ========================================================

        elif suffix == "docx":

            text = extract_docx_text(
                temp_path
            )

            st.subheader("📄 DOCX Content")

            st.text_area(
                "Extracted DOCX Text",
                text[:5000],
                height=300
            )

            with st.status("Analyzing DOCX…", expanded=True) as status:
                status.write("Mechanistic trace…")
                report = mechanistic_text(
                    text,
                    sentiment_model,
                    zero_shot_model,
                    embedding_model,
                    user_prompt,
                    use_shap=use_shap and not fast_mode,
                    fast_mode=fast_mode,
                )
                report.extend(analyze_docx_structure(text[:5000]))
                report.extend(
                    analyze_discourse(
                        text,
                        user_prompt,
                        "docx",
                        embedding_model,
                        zero_shot_model,
                    )
                )
                status.write("Generating narrative…")
            run_mechanistic_narrative(report, user_prompt)


        # ========================================================
        # PPTX
        # ========================================================

        elif suffix == "pptx":

            text = extract_pptx_text(
                temp_path
            )

            st.subheader("📄 PPTX Content")

            st.text_area(
                "Extracted PPTX Text",
                text[:5000],
                height=300
            )

            slide_blocks = [
                b for b in text.split("\n") if b.strip()
            ]
            with st.status("Analyzing PPTX…", expanded=True) as status:
                status.write("Mechanistic trace…")
                report = mechanistic_text(
                    text,
                    sentiment_model,
                    zero_shot_model,
                    embedding_model,
                    user_prompt,
                    use_shap=use_shap and not fast_mode,
                    fast_mode=fast_mode,
                )
                report.extend(
                    analyze_pptx_structure(text, slide_texts=slide_blocks)
                )
                report.extend(
                    analyze_discourse(
                        text,
                        user_prompt,
                        "pptx",
                        embedding_model,
                        zero_shot_model,
                    )
                )
                report.extend(
                    analyze_pptx_discourse(
                        slide_blocks,
                        embedding_model,
                        user_prompt,
                        zero_shot_model,
                    )
                )
                status.write("Generating narrative…")
            run_mechanistic_narrative(report, user_prompt)

        # ========================================================
        # IMAGE
        # ========================================================

        elif suffix in [

            "png",
            "jpg",
            "jpeg"

        ]:

            image = Image.open(
                temp_path
            )

            st.image(
                image,
                caption="Uploaded Image",
                use_container_width=True
            )

            explanation = explain_image(
                image
            )

            clip_m, clip_p = (get_clip() if use_clip else (None, None))
            report = mechanistic_image(
                image,
                image_model,
                imagenet_labels,
                embedding_model,
                user_prompt,
                explanation.get("Top Predictions"),
                use_clip=use_clip,
                clip_model=clip_m,
                clip_processor=clip_p,
            )
            report.artifacts["image_pil"] = image
            report.input_semantics.append(
                f"Surface type attributed: {explanation.get('Detected Image Type', 'unknown')}."
            )
            run_mechanistic_narrative(
                report,
                user_prompt,
                image=image,
                title="🧠 Visual Mechanistic Inference Narrative",
            )

        # ========================================================
        # AUDIO
        # ========================================================

        elif suffix in [

            "mp3",
            "wav"

        ]:

            st.audio(file_bytes)

            explanation = explain_audio(temp_path)

            st.subheader(
                "🎵 Audio Explainability"
            )

            transcript = explanation.get(
                "Transcript",
                ""
            )

            # ====================================================
            # HANDLE FAILED TRANSCRIPTION
            # ====================================================

            if (
                transcript.strip() == ""
                or transcript.lower().startswith(
                    "could not transcribe audio"
                )
            ):

                st.warning(
                    "Could not extract speech transcript from audio."
                )

                human_explanation = """
    The system could not clearly transcribe the uploaded audio.

    This usually happens when:
    - the speech is unclear
    - the voice is highly synthetic
    - background noise exists
    - the audio format is difficult to process
    - or speech recognition failed.

    Because no reliable transcript was extracted,
    semantic prompt-to-output reasoning cannot be performed accurately.

    Try:
    - clearer audio
    - WAV format
    - slower narration
    - or Whisper-based transcription.
    """

            else:

                y, sr_rate = librosa.load(temp_path)
                rms = librosa.feature.rms(y=y)[0]
                report = mechanistic_audio(
                    y,
                    sr_rate,
                    transcript,
                    explanation.get("Duration", 0.0),
                )
                report = enrich_audio_report(
                    report,
                    y,
                    sr_rate,
                    transcript,
                    explanation.get("Duration", 0.0),
                    user_prompt=user_prompt,
                    whisper_segments=explanation.get("Whisper Segments"),
                    zero_shot_model=zero_shot_model,
                    rms=rms,
                )
                display_mechanistic_report(report)
                human_explanation = generate_human_explanation(
                    user_prompt,
                    report_to_context(report),
                )

            st.subheader(
                "🧠 Audio Mechanistic Inference Narrative"
            )

            st.markdown(human_explanation)

        # ========================================================
        # VIDEO
        # ========================================================

        elif suffix in [

            "mp4",
            "mov",
            "avi"

        ]:

            st.video(file_bytes)

            explanation = explain_video(temp_path)

            st.subheader("🎬 Video Explainability")

            report = mechanistic_video(
                explanation.get("motion_scores", []),
                explanation.get("Duration", 0.0),
                explanation.get("FPS", 24.0),
                explanation.get("Transcript", ""),
                user_prompt=user_prompt,
                zero_shot_model=zero_shot_model,
            )
            report = enrich_video_temporal(
                report,
                explanation.get("motion_scores", []),
                explanation.get("Duration", 0.0),
                explanation.get("Transcript", ""),
                user_prompt,
                whisper_segments=explanation.get("Whisper Segments"),
                zero_shot_model=zero_shot_model,
                fps=explanation.get("FPS", 24.0),
            )
            run_mechanistic_narrative(
                report,
                user_prompt,
                title="🧠 Video Mechanistic Inference Narrative",
            )

        # ========================================================
        # UNSUPPORTED
        # ========================================================

        else:

            st.error(
                "Unsupported file type."
            )

    finally:
        safe_unlink(temp_path)
